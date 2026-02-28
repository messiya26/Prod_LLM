"""Second pass: fix remaining accent issues"""
from pptx import Presentation
import sys

FIXES2 = {
    "qualite": "qualit\u00e9",
    "difference": "diff\u00e9rence",
    "quete": "qu\u00eate",
    "chretien": "chr\u00e9tien",
    "chretiens": "chr\u00e9tiens",
    "Chretiens": "Chr\u00e9tiens",
    "eglise": "\u00e9glise",
    "carriere": "carri\u00e8re",
    "connectee": "connect\u00e9e",
    "enregistre": "enregistr\u00e9",
    " a ": " \u00e0 ",
    " a travers": " \u00e0 travers",
    " a vie": " \u00e0 vie",
    " a jour": " \u00e0 jour",
    " a l'": " \u00e0 l'",
    " a votre": " \u00e0 votre",
    " a la ": " \u00e0 la ",
    "video ": "vid\u00e9o ",
    "video,": "vid\u00e9o,",
    "Videos": "Vid\u00e9os",
    "videos": "vid\u00e9os",
    "priere": "pri\u00e8re",
    "temoignages": "t\u00e9moignages",
    "Temoignages": "T\u00e9moignages",
    "temoignage": "t\u00e9moignage",
    "defis": "d\u00e9fis",
    "Defis": "D\u00e9fis",
    "reel": "r\u00e9el",
    "acces": "acc\u00e8s",
    "Acces": "Acc\u00e8s",
    "illimite": "illimit\u00e9",
    "francophone": "francophone",
    "Francais": "Fran\u00e7ais",
    "detail": "d\u00e9tail",
    "reflete": "refl\u00e8te",
    "completes": "compl\u00e9t\u00e9s",
    "complet ": "complet ",
    "completees": "compl\u00e9t\u00e9es",
    "semantiques": "s\u00e9mantiques",
    "accessibilite": "accessibilit\u00e9",
    "Accessibilite": "Accessibilit\u00e9",
    "Viralite": "Viralit\u00e9",
    "verification": "v\u00e9rification",
    "Verification": "V\u00e9rification",
    "previsions": "pr\u00e9visions",
    "genereux": "g\u00e9n\u00e9reux",
    "Etape": "\u00c9tape",
    "etape": "\u00e9tape",
    "Actualites": "Actualit\u00e9s",
    "cles": "cl\u00e9s",
    "Cles": "Cl\u00e9s",
    "CLES": "CL\u00c9S",
    "gratuites": "gratuites",
    "creant": "cr\u00e9ant",
    "completes,": "compl\u00e9t\u00e9s,",
}

def fix_text(text):
    if not text:
        return text
    for wrong, right in sorted(FIXES2.items(), key=lambda x: -len(x[0])):
        if wrong in text and right != wrong:
            text = text.replace(wrong, right)
    return text

def process(input_path, output_path):
    prs = Presentation(input_path)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        old = run.text
                        new = fix_text(old)
                        if new != old:
                            run.text = new
                            count += 1
    prs.save(output_path)
    print(f"Pass 2: Fixed {count} more text runs. Saved: {output_path}")

if __name__ == "__main__":
    process(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else sys.argv[1])
