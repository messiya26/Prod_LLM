"""
Lord Lombo Academie - Presentation V3
Style identique a MAG VOYAGE, couleurs Lord Lombo
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

BG = RGBColor(0x0D, 0x22, 0x40)       # Navy fond
GOLD = RGBColor(0xD4, 0xA4, 0x38)      # Or principal
GOLD_LIGHT = RGBColor(0xF0, 0xC7, 0x5E) # Or clair
NAVY_DEEP = RGBColor(0x0A, 0x16, 0x28)  # Plus sombre
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CREAM = RGBColor(0xFF, 0xF8, 0xE7)
CARD_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_LIGHT = RGBColor(0xFE, 0xF9, 0xED) # Creme clair
TEXT_DARK = RGBColor(0x0D, 0x22, 0x40)
TEXT_MUTED = RGBColor(0x64, 0x74, 0x8B)
TEAL_DARK = RGBColor(0x8B, 0x69, 0x14)  # Or fonce footer
FOOTER_BG = RGBColor(0x1B, 0x3A, 0x6B)  # Bleu royal footer
BLUE_ROYAL = RGBColor(0x1B, 0x3A, 0x6B)

W_EMU = Emu(12192000)
H_EMU = Emu(6858000)
W_IN = 13.333
H_IN = 7.5

prs = Presentation()
prs.slide_width = W_EMU
prs.slide_height = H_EMU
blank_layout = prs.slide_layouts[6]

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "assets")
LOGO = os.path.join(ASSETS, "logo-client.jpeg")

def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, diameter, fill_color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(diameter), Inches(diameter))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if alpha is not None:
        shape.fill.fore_color.brightness = 0
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=14, color=WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.text_frame.word_wrap = True
    tf = txBox.text_frame
    tf.paragraphs[0].text = ""

    lines = text.split("\n")
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
    tf.auto_size = None
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    # Reduce corner radius
    shape.adjustments[0] = 0.05
    return shape

def decor_circles(slide, positions=None):
    """Add decorative circles like MAG VOYAGE"""
    if positions is None:
        positions = [
            (-2, -2, 6, GOLD, 0.08),
            (-0.7, -1, 4, GOLD, 0.06),
            (8, 3, 5.5, GOLD, 0.06),
            (9.5, 4, 3.5, GOLD, 0.08),
            (10.5, -0.8, 2.8, GOLD_LIGHT, 0.05),
        ]
    for x, y, d, color, _ in positions:
        c = add_circle(slide, x, y, d, color)
        from pptx.oxml.ns import qn
        spPr = c._element.spPr
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is None:
            fill_elem = spPr.find(qn('a:ln'))
        # Set transparency via alpha
        srgb = spPr.findall('.//' + qn('a:srgbClr'))
        for s in srgb:
            alpha_elem = s.find(qn('a:alpha'))
            if alpha_elem is None:
                from lxml import etree
                alpha_elem = etree.SubElement(s, qn('a:alpha'))
            alpha_elem.set('val', '8000')  # 8% opacity

def top_bar(slide):
    add_shape(slide, 0, 0, W_IN, 0.08, GOLD)

def footer_bar(slide, text="LORD LOMBO ACADEMIE  |  Proposition de Solution Digitale  |  Confidentiel"):
    add_shape(slide, 0, H_IN - 0.4, W_IN, 0.4, FOOTER_BG)
    add_text(slide, 0, H_IN - 0.38, W_IN, 0.35, text, font_size=9, color=RGBColor(0xB0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)

def section_slide(slide, number, title, subtitle):
    set_bg(slide, BG)
    decor_circles(slide)
    top_bar(slide)
    footer_bar(slide)
    # Number big
    add_text(slide, 1, 1.8, 2, 1, number, font_size=72, color=GOLD, bold=True, font_name="Georgia")
    # Separator
    add_shape(slide, 1, 3.0, 1.5, 0.05, GOLD)
    # Title
    add_text(slide, 1, 3.3, 10, 1.2, title, font_size=30, color=WHITE, bold=True, font_name="Georgia")
    # Subtitle
    add_text(slide, 1, 4.5, 10, 0.8, subtitle, font_size=16, color=GOLD_LIGHT, font_name="Georgia")

def sommaire_row(slide, num, text_str, y, is_even):
    bg_color = CARD_WHITE if not is_even else CARD_LIGHT
    add_rounded_rect(slide, 0.75, y, 10.8, 0.43, bg_color)
    add_shape(slide, 0.75, y, 0.55, 0.43, GOLD)
    add_text(slide, 0.75, y + 0.04, 0.55, 0.35, num, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 1.45, y + 0.04, 9, 0.35, text_str, font_size=14, color=TEXT_DARK)

# ================================================================
# SLIDE 1 - COUVERTURE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl)
top_bar(sl)

sl.shapes.add_picture(LOGO, Inches(1.1), Inches(0.8), Inches(3.5), Inches(2.5))

add_text(sl, 1.1, 1.1, 4.5, 1, "LORD LOMBO", font_size=54, color=GOLD, bold=True, font_name="Georgia")
add_text(sl, 3.5, 1.1, 5, 1, "ACADEMIE", font_size=54, color=WHITE, bold=True, font_name="Georgia")
add_text(sl, 1.1, 2.1, 7, 0.5, "PLATEFORME DE FORMATION EN LIGNE", font_size=13, color=GOLD_LIGHT)
add_shape(sl, 1.1, 2.75, 1.8, 0.04, GOLD)

add_text(sl, 1.1, 3.1, 9, 1.2,
         "DOCUMENT DE COMPREHENSION DES BESOINS\n& PROPOSITION DE SOLUTION DIGITALE",
         font_size=24, color=WHITE, bold=True, font_name="Georgia")
# Second line gold
txBox = sl.shapes[-1]
txBox.text_frame.paragraphs[1].font.color.rgb = GOLD

add_rounded_rect(sl, 1.1, 4.7, 4.2, 1, RGBColor(0x1E, 0x29, 0x3B))
add_text(sl, 1.25, 4.8, 3.9, 0.4, "Prepare par :", font_size=10, color=TEXT_MUTED)
add_text(sl, 1.25, 5.1, 3.9, 0.5, "Messiya Group - Division Digitale", font_size=12, color=GOLD_LIGHT, bold=True)

add_shape(sl, 0, H_IN - 0.4, W_IN, 0.4, FOOTER_BG)
add_text(sl, 3.5, H_IN - 0.38, 5, 0.35, "CONFIDENTIEL", font_size=10, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# ================================================================
# SLIDE 2 - SOMMAIRE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [
    (-1, -1.2, 3.5, GOLD, 0.08),
    (10, 3.5, 3.5, GOLD, 0.06),
    (11, -0.6, 2.2, GOLD_LIGHT, 0.05),
])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 5, 0.7, "SOMMAIRE", font_size=34, color=GOLD, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.95, 1.2, 0.05, GOLD)

items = [
    ("01", u"Contexte & Compr\u00e9hension du Besoin"),
    ("02", u"Vision & Objectifs Strat\u00e9giques"),
    ("03", u"Valeurs Ajout\u00e9es de Notre Solution"),
    ("04", u"Architecture Technique de la Solution"),
    ("05", u"Modules Fonctionnels de la Plateforme"),
    ("06", u"CRM Int\u00e9gr\u00e9 & Gestion Relation Client"),
    ("07", u"M\u00e9thodes P\u00e9dagogiques & Organisation"),
    ("08", u"\u00c9valuation, Satisfaction & Suivi"),
    ("09", u"Fonctionnalit\u00e9s Avanc\u00e9es"),
    ("10", u"S\u00e9curit\u00e9 & Performance"),
    ("11", u"Identit\u00e9 Visuelle & Charte Graphique"),
    ("12", u"Stack Technique & Technologies"),
    ("13", u"Strat\u00e9gie SEO & R\u00e9f\u00e9rencement"),
    ("14", u"R\u00e9seaux Sociaux & Marketing Digital"),
    ("15", u"Parcours Utilisateur & Exp\u00e9rience (UX)"),
    ("16", u"Contact & Devis Personnalis\u00e9"),
    ("17", u"\u00c9volutions Futures & Roadmap"),
    ("18", u"Valeur Ajout\u00e9e Messiya Group"),
]
for i, (num, txt) in enumerate(items):
    sommaire_row(sl, num, txt, 1.25 + i * 0.335, i % 2 == 1)

# ================================================================
# SLIDE 3 - SECTION: CONTEXTE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "01", u"CONTEXTE & COMPR\u00c9HENSION\nDU BESOIN",
              u"Comprendre la vision de Lord Lombo pour une acad\u00e9mie de transformation spirituelle")

# ================================================================
# SLIDE 4 - CONTEXTE DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1.5, -1.5, 4.5, GOLD, 0.06), (10, 4, 4, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 8, 0.6, "CONTEXTE DU PROJET", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# Context box
add_rounded_rect(sl, 0.5, 1.2, 7, 2.2, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.7, 1.3, 2.2, 0.4, "LORD LOMBO ACADEMIE", font_size=14, color=GOLD, bold=True)
add_text(sl, 0.7, 1.7, 6.5, 1.5,
         "est une plateforme de formation en ligne dediee a la transformation spirituelle, "
         "au leadership et au developpement personnel. Elle accompagne des individus dans leur "
         "croissance interieure, renforce leur influence en tant que leaders et leur permet "
         "d'impacter positivement leur environnement.\n\n"
         "L'academie offrira des formations certifiantes, des programmes de coaching, "
         "des sessions interactives en direct et un suivi personnalise de chaque apprenant.",
         font_size=11, color=RGBColor(0xB0, 0xC0, 0xD0))

# 4 numbered cards
cards = [
    ("1", "E-Learning Premium", "Formations video HD, quiz, documents PDF, acces illimite 24h/24"),
    ("2", "Coaching Live", "Sessions Zoom/Meet integrees, interaction temps reel, replays"),
    ("3", "Communaute", "Forum, defis spirituels, temoignages, espace de priere collectif"),
    ("4", "Certification", "Certificats PDF automatiques avec QR code de verification"),
]
for i, (num, title, desc) in enumerate(cards):
    x = 0.5 + i * 3.1
    add_rounded_rect(sl, x, 3.7, 2.9, 1.6, CARD_WHITE)
    add_shape(sl, x, 3.7, 0.45, 0.45, GOLD)
    add_text(sl, x, 3.73, 0.45, 0.4, num, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + 0.55, 3.75, 2.2, 0.4, title, font_size=12, color=TEXT_DARK, bold=True)
    add_text(sl, x + 0.15, 4.25, 2.6, 0.9, desc, font_size=10, color=TEXT_MUTED)

# Right side: Cible
add_rounded_rect(sl, 7.7, 1.2, 5.2, 4.1, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 7.9, 1.35, 3, 0.4, "PUBLIC CIBLE", font_size=14, color=GOLD, bold=True)
targets = [
    "Leaders en quete de croissance interieure",
    "Chretiens desireux de se former en profondeur",
    "Jeunes leaders en developpement",
    "Communaute francophone mondiale",
    "Pasteurs et responsables d'eglise",
    "Professionnels en reconversion de carriere",
    "Diaspora africaine connectee",
    "Entrepreneurs chretiens",
]
for i, t in enumerate(targets):
    add_text(sl, 7.9, 1.8 + i * 0.4, 4.8, 0.35, f"  {t}", font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# Exigences bottom
add_rounded_rect(sl, 7.7, 5.5, 5.2, 1.1, GOLD)
add_text(sl, 7.9, 5.55, 4.8, 0.3, "EXIGENCES CLES DU CLIENT", font_size=11, color=TEXT_DARK, bold=True)
add_text(sl, 7.9, 5.85, 4.8, 0.7, "Design haut de gamme  |  Securite bancaire  |  Multi-paiement\nMobile-first  |  Certificats auto  |  Coaching live & enregistre", font_size=10, color=RGBColor(0x1A, 0x2D, 0x50))

# ================================================================
# SLIDE 5 - SECTION: VISION
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "02", "VISION & OBJECTIFS\nSTRATEGIQUES",
              "Former, Transformer, Impacter - La triple mission de Lord Lombo Academie")

# ================================================================
# SLIDE 6 - VISION DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3.5, GOLD, 0.06), (9.5, 4.5, 4, GOLD, 0.06), (10.5, -0.5, 2.5, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 8, 0.6, "VISION & OBJECTIFS", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# 3 big cards
pillars = [
    ("FORMER", "Offrir des parcours de formation structurees et certifiantes en croissance spirituelle, leadership chretien et influence positive"),
    ("TRANSFORMER", "Accompagner chaque apprenant dans une transformation profonde et durable a travers un suivi personnalise, du coaching et des defis pratiques"),
    ("IMPACTER", "Creer une communaute de leaders influents qui impactent positivement leur environnement en Afrique et dans la diaspora francophone"),
]
for i, (title, desc) in enumerate(pillars):
    x = 0.5 + i * 4.1
    add_rounded_rect(sl, x, 1.2, 3.9, 2.8, CARD_WHITE)
    add_shape(sl, x, 1.2, 3.9, 0.5, GOLD)
    add_text(sl, x, 1.23, 3.9, 0.45, title, font_size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font_name="Georgia")
    add_text(sl, x + 0.2, 1.85, 3.5, 2, desc, font_size=12, color=TEXT_MUTED)

# Vision quote
add_rounded_rect(sl, 0.5, 4.3, 12.3, 1.5, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.8, 4.5, 11.7, 1,
         "\"Devenir LA reference dans l'espace francophone pour la formation spirituelle,\n"
         "le leadership et le developpement personnel, en combinant technologie de pointe\n"
         "et accompagnement humain de qualite.\"",
         font_size=16, color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")
add_shape(sl, 5.5, 5.55, 2.3, 0.04, GOLD)
add_text(sl, 3, 5.65, 7, 0.3, "Lord Lombo Academie - Vision 2026", font_size=10, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 7 - SECTION: VALEURS AJOUTEES
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "03", "VALEURS AJOUTEES\nDE NOTRE SOLUTION",
              "Ce que nous apportons en plus pour faire la difference")

# ================================================================
# SLIDE 8 - VALEURS AJOUTEES DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1.5, -1, 4, GOLD, 0.06), (10, 3.5, 4.5, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "VALEUR AJOUTEE MESSIYA GROUP", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)
add_text(sl, 0.75, 0.95, 10, 0.4, "Ce que nous apportons en plus pour faire la difference", font_size=13, color=GOLD_LIGHT)

values = [
    ("Experience Premium", "Interface haut de gamme, animations fluides Framer Motion, micro-interactions. Chaque detail reflete l'excellence de la vision de Lord Lombo."),
    ("Intelligence Artificielle", "Recommandations personnalisees de formations, sous-titres auto (Whisper AI), chatbot assistant 24/7 pour les apprenants."),
    ("Multi-langue & Accessibilite", "Support Francais, Anglais, Lingala + accessibilite WCAG 2.1 AA pour toucher un maximum d'apprenants dans le monde."),
    ("Programme de Fidelite", "Points de fidelite, badges gamifies, recompenses pour les defis completes, parrainage avec avantages automatiques."),
    ("Tableau de Bord Admin", "Back-office complet : gestion apprenants, formations, paiements, analytics en temps reel, rapports de satisfaction."),
    ("Maintenance & Evolution", "Support technique continu, mises a jour securite, evolution des fonctionnalites, SLA 99.9% disponibilite garantie."),
]
for i, (title, desc) in enumerate(values):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.5 + row * 1.7
    add_rounded_rect(sl, x, y, 6, 1.5, RGBColor(0x14, 0x28, 0x42))
    add_shape(sl, x, y, 0.08, 1.5, GOLD)
    add_text(sl, x + 0.25, y + 0.1, 0.4, 0.4, "+", font_size=22, color=GOLD, bold=True)
    add_text(sl, x + 0.7, y + 0.12, 5, 0.35, title, font_size=14, color=WHITE, bold=True)
    add_text(sl, x + 0.7, y + 0.5, 5, 0.85, desc, font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# ================================================================
# SLIDE 9 - SECTION: ARCHITECTURE TECHNIQUE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "04", "ARCHITECTURE TECHNIQUE\nDE LA SOLUTION",
              "Une infrastructure moderne, scalable et performante pour une acad\u00e9mie de r\u00e9f\u00e9rence")

# ================================================================
# SLIDE 10 - ARCHITECTURE TECHNIQUE DETAIL (comme MAG slide 10)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3.5, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "ARCHITECTURE DE LA SOLUTION", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# 3 couches comme MAG VOYAGE
arch_layers = [
    ("COUCHE PR\u00c9SENTATION (FRONTEND)", [
        ("Next.js 14+", "SSR / SSG"),
        ("React 18", "Composants UI"),
        ("Tailwind CSS", "Design System"),
        ("Framer Motion", "Animations"),
        ("PWA", "Mode Offline"),
    ]),
    ("COUCHE M\u00c9TIER (API & BACKEND)", [
        ("NestJS", "REST / GraphQL"),
        ("Auth.js", "JWT / 2FA"),
        ("Prisma ORM", "Acc\u00e8s donn\u00e9es"),
        ("Bull MQ", "Jobs async"),
        ("WebSocket", "Temps r\u00e9el"),
    ]),
    ("DONN\u00c9ES & STOCKAGE", [
        ("PostgreSQL 16", "Base relationnelle"),
        ("Redis 7", "Cache & sessions"),
        ("Mux / CF Stream", "Vid\u00e9o CDN"),
    ]),
]

y_pos = 1.15
for layer_name, techs in arch_layers:
    add_rounded_rect(sl, 0.5, y_pos, 8, 0.4, GOLD)
    add_text(sl, 0.6, y_pos + 0.03, 7.5, 0.35, layer_name, font_size=11, color=TEXT_DARK, bold=True)
    y_pos += 0.5
    for j, (tech, role) in enumerate(techs):
        x = 0.5 + j * 1.6
        w = 1.45
        add_rounded_rect(sl, x, y_pos, w, 0.6, CARD_WHITE)
        add_text(sl, x + 0.05, y_pos + 0.03, w - 0.1, 0.25, tech, font_size=9, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl, x + 0.05, y_pos + 0.3, w - 0.1, 0.25, role, font_size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    y_pos += 0.8
    if layer_name != "DONN\u00c9ES & STOCKAGE":
        add_text(sl, 4, y_pos - 0.15, 1, 0.2, "\u25bc", font_size=14, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

add_rounded_rect(sl, 0.5, y_pos + 0.1, 8, 0.4, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.6, y_pos + 0.13, 7.5, 0.35,
    "SERVICES EXTERNES : Stripe / PayPal / Mobile Money  |  Resend  |  Zoom SDK  |  OneSignal",
    font_size=9, color=RGBColor(0xB0, 0xC0, 0xD0))

# Right side
add_rounded_rect(sl, 8.8, 1.15, 4.2, 3.4, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 9, 1.25, 3.8, 0.35, "INFRASTRUCTURE & DEVOPS", font_size=12, color=GOLD, bold=True)
infra_arch = [
    "AWS (EC2, RDS, S3, CloudFront)",
    "Docker + GitHub Actions CI/CD",
    "Cloudflare CDN & WAF",
    "Terraform (IaC)",
    "Datadog monitoring 24/7",
    "Auto-scaling horizontal",
]
for i, item in enumerate(infra_arch):
    add_text(sl, 9, 1.7 + i * 0.36, 3.8, 0.33, f"\u2713  {item}", font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

add_rounded_rect(sl, 8.8, 4.7, 4.2, 1.8, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 9, 4.8, 3.8, 0.35, "VID\u00c9O & STREAMING", font_size=12, color=GOLD, bold=True)
vid_arch = [
    "Mux / Cloudflare Stream",
    "HLS adaptive bitrate",
    "Protection DRM contenu",
    "Sous-titres auto (Whisper AI)",
]
for i, item in enumerate(vid_arch):
    add_text(sl, 9, 5.2 + i * 0.3, 3.8, 0.28, f"\u2713  {item}", font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

add_rounded_rect(sl, 0.5, 6.2, 12.3, 0.5, GOLD)
add_text(sl, 0.5, 6.25, 12.3, 0.4,
    "Vercel (CDN Global)  |  Supabase (BaaS)  |  Cloudflare (DDoS)  |  SSL/TLS  |  CI/CD GitHub Actions  |  Analytics GA4",
    font_size=10, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 11 - SECTION: MODULES
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "05", "MODULES FONCTIONNELS\nDE LA PLATEFORME",
              "L'ensemble des pages et espaces qui composent l'acad\u00e9mie en ligne")

# ================================================================
# SLIDE 12 - MODULES DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1.5, 3.5, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 8, 0.6, "MODULES DE LA PLATEFORME", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

modules = [
    ("01", "Accueil Inspirant", "Hero video Lord Lombo, vision, CTA transformation, temoignages, formations phares"),
    ("02", "A Propos & Mission", "Presentation Lord Lombo, histoire academie, equipe pedagogique, fondements spirituels"),
    ("03", "Catalogue Formations", "Filtres par theme, fiches detaillees, formateur, duree, prix, avis, prerequis"),
    ("04", "Espace Membre", "Dashboard, progression %, certificats, historique, defis pratiques, parametres"),
    ("05", "Paiement & Facturation", "Stripe, PayPal, Mobile Money, abonnements, factures PDF auto, codes promo"),
    ("06", "Sessions Live", "Integration Zoom/Meet, notifications, replays, chat, interaction temps reel"),
    ("07", "Forum & Communaute", "Echanges, temoignages, entraide, espace priere, defis collectifs, moderation"),
    ("08", "Blog & Contenu SEO", "Articles inspires, guides spirituels, temoignages, newsletters, partage social"),
    ("09", "Administration", "Gestion users, formations, paiements, stats temps reel, rapports, moderation"),
]
for i, (num, title, desc) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.2 + row * 1.65
    add_rounded_rect(sl, x, y, 3.95, 1.45, CARD_WHITE)
    add_shape(sl, x, y, 0.5, 0.4, GOLD)
    add_text(sl, x, y + 0.03, 0.5, 0.35, num, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + 0.6, y + 0.05, 3.2, 0.35, title, font_size=12, color=TEXT_DARK, bold=True)
    add_text(sl, x + 0.15, y + 0.5, 3.6, 0.85, desc, font_size=10, color=TEXT_MUTED)

# ================================================================
# SLIDE - SECTION: CRM & GESTION RELATION CLIENT
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "06", "CRM INT\u00c9GR\u00c9 & GESTION\nRELATION CLIENT",
              "Un pipeline automatis\u00e9 pour convertir les visiteurs en ambassadeurs de la vision")

# ================================================================
# SLIDE - CRM DETAIL (pipeline comme MAG)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "CRM & GESTION RELATION CLIENT", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# Pipeline visuel
pipeline_steps = ["Visiteur", "\u25b6", "Inscrit", "\u25b6", "Apprenant", "\u25b6", "Certifi\u00e9", "\u25b6", "Ambassadeur"]
pipeline_text = "   ".join(pipeline_steps)
add_rounded_rect(sl, 0.5, 1.15, 12.3, 0.55, GOLD)
add_text(sl, 0.6, 1.2, 12.1, 0.45, pipeline_text, font_size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

crm_cards = [
    ("Acquisition & Inscription", "Formulaire intelligent avec SSO (Google, Facebook), tracking source d'acquisition (UTM), segmentation automatique par int\u00e9r\u00eat (spirituel, leadership, d\u00e9veloppement)."),
    ("Engagement & Suivi", "Emails de bienvenue automatis\u00e9s, rappels d'inactivit\u00e9, notifications de nouveaux cours, push personnalis\u00e9s selon le profil et la progression."),
    ("Conversion & Paiement", "Relances panier abandonn\u00e9, offres personnalis\u00e9es, codes promo cibl\u00e9s, upsell automatique vers formations premium ou coaching."),
    ("Fid\u00e9lisation & R\u00e9tention", "Programme de fid\u00e9lit\u00e9 (points, badges), anniversaires, offres de renouvellement d'abonnement, NPS automatique."),
    ("Parrainage & Ambassadeur", "Syst\u00e8me d'affiliation int\u00e9gr\u00e9, lien de parrainage unique, r\u00e9compenses automatiques, suivi des conversions par ambassadeur."),
    ("Analytics & Pilotage", "Tableau de bord CRM : entonnoir de conversion, LTV, taux de r\u00e9tention, cohortes, pr\u00e9visions de revenus, alertes automatiques."),
]
for i, (title, desc) in enumerate(crm_cards):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.95 + row * 1.55
    add_rounded_rect(sl, x, y, 6, 1.35, CARD_WHITE)
    add_shape(sl, x, y, 0.08, 1.35, GOLD)
    add_text(sl, x + 0.25, y + 0.1, 5.5, 0.35, title, font_size=13, color=TEXT_DARK, bold=True)
    add_text(sl, x + 0.25, y + 0.45, 5.5, 0.8, desc, font_size=10, color=TEXT_MUTED)

# Bottom
add_rounded_rect(sl, 0.5, 6.2, 12.3, 0.45, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.5, 6.25, 12.3, 0.35,
    "Int\u00e9grations : Mailchimp  |  Resend  |  WhatsApp Business  |  Google Analytics 4  |  Facebook Pixel  |  OneSignal",
    font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE - SECTION: METHODES PEDAGOGIQUES
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "07", u"M\u00c9THODES P\u00c9DAGOGIQUES\n& ORGANISATION DES COURS",
              u"Un mod\u00e8le hybride innovant combinant autonomie et accompagnement humain")

# ================================================================
# SLIDE 12 - METHODES DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (9.5, 4, 4.5, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "ORGANISATION DES COURS", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

methods = [
    ("COURS ENREGISTRES", [
        "Videos HD pre-enregistrees 24h/24",
        "Documents PDF complementaires",
        "Fichiers audio (podcasts spirituels)",
        "Quiz interactifs par module",
        "Acces illimite a vie",
        "Suivi automatique de progression",
        "Sous-titres multilingues (AI)",
    ]),
    ("SESSIONS EN DIRECT", [
        "Integration Zoom / Google Meet",
        "Lien dans l'espace membre",
        "Notifications & rappels auto",
        "Interaction temps reel (Q&A)",
        "Moments de priere collective",
        "Coaching en direct",
        "Replays disponibles pour absents",
    ]),
    ("PARCOURS PERSONNALISE", [
        "Mix videos + sessions live",
        "Coaching individuel premium",
        "Coaching de groupe",
        "Defis pratiques hebdomadaires",
        "Exercices d'application concrets",
        "Mentorat personnalise",
        "Suivi par formateur dedie",
    ]),
]
for i, (title, items_list) in enumerate(methods):
    x = 0.5 + i * 4.15
    add_rounded_rect(sl, x, 1.15, 3.95, 4.2, CARD_WHITE)
    add_shape(sl, x, 1.15, 3.95, 0.5, GOLD)
    add_text(sl, x, 1.18, 3.95, 0.45, title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items_list):
        add_text(sl, x + 0.2, 1.75 + j * 0.45, 3.5, 0.4, f"  {item}", font_size=10, color=TEXT_MUTED)

# Bottom bar
add_rounded_rect(sl, 0.5, 5.55, 12.3, 0.9, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.8, 5.6, 3, 0.35, "MODELE HYBRIDE INNOVANT", font_size=12, color=GOLD, bold=True)
add_text(sl, 0.8, 5.95, 11.7, 0.4,
         "La combinaison de ces 3 approches garantit une experience complete, flexible et adaptee a chaque profil. "
         "L'apprenant choisit son rythme tout en beneficiant d'un accompagnement humain de qualite.",
         font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# ================================================================
# SLIDE 13 - SECTION: EVALUATION
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "06", "EVALUATION, SATISFACTION\n& SUIVI DE PROGRESSION",
              "Mesurer l'impact des formations et accompagner chaque apprenant vers l'excellence")

# ================================================================
# SLIDE 14 - EVALUATION DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3.5, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 8, 0.6, "EVALUATION & SATISFACTION", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# Pipeline flow
steps = ["Inscription", ">", "Formation", ">", "Evaluation", ">", "Certificat", ">", "Temoignage", ">", "Ambassadeur"]
flow_text = "  ".join(steps)
add_rounded_rect(sl, 0.5, 1.15, 12.3, 0.55, GOLD)
add_text(sl, 0.6, 1.2, 12.1, 0.45, flow_text, font_size=13, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

eval_cards = [
    ("Questionnaires Auto", "Apres chaque module et formation, questionnaire de satisfaction envoye automatiquement a l'apprenant."),
    ("Notation Etoiles", "Chaque formation, formateur et module peut etre note individuellement (1 a 5 etoiles) avec commentaires."),
    ("Rapports Satisfaction", "Dashboard admin dedie avec taux de satisfaction, NPS, tendances, alertes automatiques."),
    ("Suivi Progression", "Barre de progression visuelle (%), badges gamifies, historique complet, rappels automatiques."),
    ("KPI Detailles", "Taux completion, score moyen, retention abonnes, inscriptions/mois, revenus par formation."),
    ("Export & Partage", "Export resultats PDF/Excel, certificats partageables LinkedIn, temoignages indexes sur le site."),
]
for i, (title, desc) in enumerate(eval_cards):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.95 + row * 2.15
    add_rounded_rect(sl, x, y, 3.95, 1.95, CARD_WHITE)
    add_shape(sl, x, y, 0.08, 1.95, GOLD)
    add_text(sl, x + 0.25, y + 0.15, 3.5, 0.35, title, font_size=13, color=TEXT_DARK, bold=True)
    add_text(sl, x + 0.25, y + 0.55, 3.5, 1.2, desc, font_size=11, color=TEXT_MUTED)

# ================================================================
# SLIDE 15 - SECTION: FONCTIONNALITES
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "07", "FONCTIONNALITES\nAVANCEES",
              "Les outils qui font la difference pour une academie complete et performante")

# ================================================================
# SLIDE 16 - FONCTIONNALITES DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1.5, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "FONCTIONNALITES AVANCEES", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

features = [
    ("Responsive & PWA", "Mobile-first, installable, offline"),
    ("Certificats PDF", "QR code, verification en ligne"),
    ("Notifications Smart", "Email, push, SMS, rappels auto"),
    ("Forum Communaute", "Echanges, entraide, temoignages"),
    ("Calendrier Events", "Sessions live, webinaires, conferences"),
    ("Dashboard Analytics", "KPI, rapports, stats temps reel"),
    ("Evaluation & Avis", "Etoiles, NPS, questionnaires"),
    ("Multi-Paiement", "Carte, PayPal, Mobile Money, promo"),
    ("Securite Renforcee", "SSL, 2FA, RGPD, WAF, DDoS"),
    ("Defis Spirituels", "Exercices pratiques quotidiens"),
    ("Streaming Video", "HLS adaptatif, DRM, multi-qualite"),
    ("Email Marketing", "Sequences auto, templates, relances"),
]
for i, (title, desc) in enumerate(features):
    col = i % 4
    row = i // 4
    x = 0.3 + col * 3.2
    y = 1.15 + row * 1.65
    add_rounded_rect(sl, x, y, 3, 1.45, CARD_WHITE)
    add_shape(sl, x, y, 3, 0.04, GOLD)
    add_text(sl, x + 0.15, y + 0.15, 2.7, 0.4, title, font_size=12, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + 0.15, y + 0.55, 2.7, 0.7, desc, font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Integration bar
add_text(sl, 0.3, 6.2, 12.7, 0.35,
         "+ Integration WhatsApp Business  |  Google Analytics 4  |  Zoom SDK  |  Stripe Webhooks  |  Mailchimp",
         font_size=10, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 17 - SECTION: SECURITE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "08", "SECURITE &\nPERFORMANCE",
              "Des standards de securite maximaux pour proteger les donnees et les transactions")

# ================================================================
# SLIDE 18 - SECURITE DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3.5, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 8, 0.6, "SECURITE & PAIEMENTS", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

sec_items = [
    ("HTTPS / SSL", "Certificat SSL, chiffrement bout en bout sur tout le site"),
    ("Authentification", "JWT + 2FA obligatoire pour admin, optionnel pour membres"),
    ("Protection DDoS", "Cloudflare WAF en frontal, rate limiting intelligent"),
    ("Anti-Injection", "Protection SQL injection, XSS, CSRF sur toutes les API"),
    ("Sauvegardes", "Backups quotidiens automatiques, retention 30 jours, 3 sites"),
    ("RGPD", "Consentement cookies, politique confidentialite, droit a l'oubli"),
    ("Rate Limiting", "Anti brute force, protection abus API, throttling intelligent"),
    ("Audit Securite", "Scan vulnerabilites trimestriel, tests penetration avant prod"),
]
for i, (title, desc) in enumerate(sec_items):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.15 + row * 1.05
    add_rounded_rect(sl, x, y, 6, 0.9, RGBColor(0x14, 0x28, 0x42))
    add_shape(sl, x, y, 0.08, 0.9, GOLD)
    add_text(sl, x + 0.25, y + 0.08, 2.5, 0.35, title, font_size=12, color=GOLD, bold=True)
    add_text(sl, x + 0.25, y + 0.42, 5.5, 0.4, desc, font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# Payment methods
add_rounded_rect(sl, 0.5, 5.45, 12.3, 0.65, GOLD)
add_text(sl, 0.5, 5.5, 12.3, 0.3, "MOYENS DE PAIEMENT SECURISES", font_size=12, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, 0.5, 5.8, 12.3, 0.3, "Stripe (3D Secure)  |  PayPal Business  |  Orange Money  |  MTN MoMo  |  Wave  |  Airtel Money  |  Flutterwave", font_size=11, color=RGBColor(0x1A, 0x2D, 0x50), align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 19 - SECTION: IDENTITE VISUELLE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "09", "IDENTITE VISUELLE &\nCHARTE GRAPHIQUE",
              "Une charte graphique premium inspiree de la spiritualite et du leadership")

# ================================================================
# SLIDE 20 - IDENTITE VISUELLE DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 3, 4.5, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "PALETTE DE COULEURS & DESIGN", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

colors_list = [
    ("#0D2240", "Bleu Nuit", "Headers, navigation", RGBColor(0x0D, 0x22, 0x40), WHITE),
    ("#1B3A6B", "Bleu Royal", "Accents principaux", RGBColor(0x1B, 0x3A, 0x6B), WHITE),
    ("#D4A438", "Or Royal", "Elements premium", RGBColor(0xD4, 0xA4, 0x38), TEXT_DARK),
    ("#F0C75E", "Or Clair", "CTA, boutons", RGBColor(0xF0, 0xC7, 0x5E), TEXT_DARK),
    ("#0A1628", "Nuit Profonde", "Fonds sombres", RGBColor(0x0A, 0x16, 0x28), WHITE),
    ("#142842", "Marine", "Cards sombres", RGBColor(0x14, 0x28, 0x42), WHITE),
    ("#FEF9ED", "Creme", "Fonds clairs", RGBColor(0xFE, 0xF9, 0xED), TEXT_DARK),
    ("#1E293B", "Ardoise", "Textes, footer", RGBColor(0x1E, 0x29, 0x3B), WHITE),
]
for i, (hex_val, name, usage, rgb, txt_color) in enumerate(colors_list):
    x = 0.5 + i * 1.55
    add_rounded_rect(sl, x, 1.15, 1.4, 1.6, rgb)
    add_text(sl, x, 1.25, 1.4, 0.35, hex_val, font_size=9, color=txt_color, align=PP_ALIGN.CENTER, bold=True)
    add_text(sl, x, 1.6, 1.4, 0.35, name, font_size=10, color=txt_color, align=PP_ALIGN.CENTER)
    add_text(sl, x, 1.95, 1.4, 0.35, usage, font_size=8, color=txt_color, align=PP_ALIGN.CENTER)

# Typographies
add_text(sl, 0.5, 3.05, 12, 0.35, "Typographies :  Georgia (Titres)  |  Inter / Calibri (Corps)  |  Georgia Italic (Citations & accents)", font_size=11, color=GOLD_LIGHT, align=PP_ALIGN.CENTER)

# Design principles
add_rounded_rect(sl, 0.5, 3.6, 6, 2.8, CARD_WHITE)
add_shape(sl, 0.5, 3.6, 6, 0.45, GOLD)
add_text(sl, 0.5, 3.63, 6, 0.4, "PRINCIPES DE DESIGN", font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
principles = [
    "Design professionnel, sobre et inspirant",
    "Interface claire et intuitive (UX premium)",
    "Integration citations & messages inspirants",
    "Responsive design mobile-first",
    "Accessibilite WCAG 2.1 AA",
    "Animations subtiles (Framer Motion)",
    "Hierarchie visuelle forte & espacement genereux",
]
for i, p in enumerate(principles):
    add_text(sl, 0.7, 4.15 + i * 0.3, 5.5, 0.28, f"  {p}", font_size=10, color=TEXT_MUTED)

# Logo
sl.shapes.add_picture(LOGO, Inches(7), Inches(3.6), Inches(3), Inches(2.3))

# Engagement
add_rounded_rect(sl, 10.2, 3.6, 2.6, 2.8, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 10.3, 3.7, 2.4, 0.35, "ENGAGEMENTS", font_size=11, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
engagements = ["Dark mode", "Micro-interactions", "Core Web Vitals", "Lighthouse > 90", "100% responsive"]
for i, e in enumerate(engagements):
    add_text(sl, 10.3, 4.15 + i * 0.42, 2.4, 0.35, e, font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 21 - SECTION: STACK TECHNIQUE
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "10", "STACK TECHNIQUE\n& TECHNOLOGIES",
              "Des technologies de pointe pour des performances exceptionnelles")

# ================================================================
# SLIDE 22 - STACK TECHNIQUE DETAIL (Architecture)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "ARCHITECTURE DE LA SOLUTION", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# 3 layer cards
layers = [
    ("COUCHE PRESENTATION (FRONTEND)", [
        ("Next.js 14+", "SSR / SSG"),
        ("React 18", "Composants UI"),
        ("Tailwind CSS", "Design System"),
        ("Framer Motion", "Animations"),
        ("PWA", "Mode Offline"),
    ]),
    ("COUCHE METIER (API & BACKEND)", [
        ("NestJS", "REST / GraphQL"),
        ("Auth.js", "JWT / 2FA"),
        ("Prisma ORM", "Acces donnees"),
        ("Bull MQ", "Jobs async"),
        ("WebSocket", "Temps reel"),
    ]),
    ("DONNEES & STOCKAGE", [
        ("PostgreSQL 16", "Base relationnelle"),
        ("Redis 7", "Cache & sessions"),
        ("Mux / Cloudflare", "Video CDN"),
    ]),
]
y_pos = 1.15
for layer_name, techs in layers:
    add_rounded_rect(sl, 0.5, y_pos, 8, 0.4, GOLD)
    add_text(sl, 0.6, y_pos + 0.03, 7.5, 0.35, layer_name, font_size=11, color=TEXT_DARK, bold=True)
    y_pos += 0.5
    for j, (tech, role) in enumerate(techs):
        x = 0.5 + j * 1.6
        w = 1.45
        add_rounded_rect(sl, x, y_pos, w, 0.6, CARD_WHITE)
        add_text(sl, x + 0.05, y_pos + 0.03, w - 0.1, 0.25, tech, font_size=9, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
        add_text(sl, x + 0.05, y_pos + 0.3, w - 0.1, 0.25, role, font_size=8, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    y_pos += 0.8
    # Arrow
    if layer_name != "DONNEES & STOCKAGE":
        add_text(sl, 4, y_pos - 0.15, 1, 0.2, "v", font_size=14, color=GOLD, align=PP_ALIGN.CENTER, bold=True)

# Services externes
add_rounded_rect(sl, 0.5, y_pos + 0.15, 8, 0.4, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.6, y_pos + 0.18, 7.5, 0.35, "SERVICES : Stripe / PayPal / MoMo  |  Resend (emails)  |  Zoom SDK  |  OneSignal (push)", font_size=9, color=RGBColor(0xB0, 0xC0, 0xD0))

# Right side: Additional tech
add_rounded_rect(sl, 8.8, 1.15, 4.2, 5.2, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 9, 1.25, 3.8, 0.35, "INFRASTRUCTURE & DEVOPS", font_size=12, color=GOLD, bold=True)
infra = [
    "AWS (EC2, RDS, S3, CloudFront)",
    "Docker + GitHub Actions CI/CD",
    "Cloudflare CDN & WAF",
    "Terraform (Infrastructure as Code)",
    "Datadog monitoring 24/7",
    "Auto-scaling horizontal",
    "Sentry (error tracking)",
    "Vercel (frontend hosting)",
]
for i, item in enumerate(infra):
    add_text(sl, 9, 1.7 + i * 0.38, 3.8, 0.35, f"  {item}", font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# Performance metrics
add_rounded_rect(sl, 8.8, 4.0, 4.2, 0.5, GOLD)
add_text(sl, 8.8, 4.03, 4.2, 0.45, "VIDEO & STREAMING", font_size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
vid_items = ["Mux / Cloudflare Stream", "HLS adaptive bitrate", "Protection DRM contenu", "Sous-titres auto (Whisper AI)", "Transcodage multi-qualite"]
for i, item in enumerate(vid_items):
    add_text(sl, 9, 4.6 + i * 0.33, 3.8, 0.3, f"  {item}", font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# Bottom perf bar
add_rounded_rect(sl, 0.5, 6.2, 12.3, 0.5, GOLD)
add_text(sl, 0.5, 6.25, 12.3, 0.4,
         "Lighthouse > 90   |   < 2s chargement   |   99.9% uptime   |   100% responsive   |   Chrome / Safari / Firefox / Edge",
         font_size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 23 - SECTION: SEO
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "11", "STRATEGIE SEO &\nREFERENCEMENT",
              "Referencement naturel et technique pour une visibilite maximale sur Google")

# ================================================================
# SLIDE 24 - SEO DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "STRATEGIE SEO & REFERENCEMENT", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

seo_cols = [
    ("SEO TECHNIQUE", [
        "URLs semantiques propres",
        "Lighthouse > 90",
        "Schema.org (Course, Org)",
        "Sitemap XML auto",
        "SSR (Next.js)",
        "Core Web Vitals",
    ]),
    ("SEO CONTENU", [
        "Blog optimise (articles longs)",
        "Pages 2000+ mots par formation",
        "Multi-langue FR/EN",
        "Rich snippets (FAQ, reviews)",
        "Internal linking strategique",
        "Contenu evergreen spirituel",
    ]),
    ("MARKETING DIGITAL", [
        "Google Analytics 4",
        "Google Search Console",
        "Facebook Pixel + TikTok Pixel",
        "Open Graph / Twitter Cards",
        "Google Tag Manager",
        "A/B testing & retargeting",
    ]),
]
for i, (title, items_list) in enumerate(seo_cols):
    x = 0.5 + i * 4.15
    add_rounded_rect(sl, x, 1.15, 3.95, 3.8, CARD_WHITE)
    add_shape(sl, x, 1.15, 3.95, 0.45, GOLD)
    add_text(sl, x, 1.18, 3.95, 0.4, title, font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items_list):
        add_text(sl, x + 0.2, 1.7 + j * 0.45, 3.5, 0.4, f"  {item}", font_size=11, color=TEXT_MUTED)

# Keywords bar
add_rounded_rect(sl, 0.5, 5.2, 12.3, 0.9, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.7, 5.25, 11.9, 0.3, "MOTS-CLES CIBLES :", font_size=11, color=GOLD, bold=True)
add_text(sl, 0.7, 5.55, 11.9, 0.5,
         "formation spirituelle en ligne | leadership chretien | academie lord lombo | developpement personnel chretien | "
         "coaching spirituel | formation leadership afrique | e-learning chretien francophone | transformation spirituelle",
         font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# ================================================================
# SLIDE 25 - SECTION: RESEAUX SOCIAUX
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "12", "RESEAUX SOCIAUX &\nMARKETING DIGITAL",
              "Strategie multi-plateforme pour maximiser la visibilite et l'engagement")

# ================================================================
# SLIDE 26 - RESEAUX SOCIAUX DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "STRATEGIE RESEAUX SOCIAUX", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

socials = [
    ("Facebook", "Communaute", "Temoignages, lives, promos, events"),
    ("Instagram", "Inspiration", "Photos, Reels, Stories behind-the-scenes"),
    ("TikTok", "Viralite", "Videos courtes, challenges spirituels"),
    ("YouTube", "Long format", "Vlogs, sermons, masterclass gratuites"),
    ("Twitter/X", "Actualites", "News, offres flash, threads inspirants"),
]
for i, (platform, category, desc) in enumerate(socials):
    x = 0.5 + i * 2.5
    add_rounded_rect(sl, x, 1.15, 2.3, 2.2, CARD_WHITE)
    add_shape(sl, x, 1.15, 2.3, 0.04, GOLD)
    add_text(sl, x, 1.3, 2.3, 0.35, platform, font_size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, 1.65, 2.3, 0.3, category, font_size=10, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_text(sl, x + 0.1, 2, 2.1, 1.2, desc, font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# Integrations
add_rounded_rect(sl, 0.5, 3.6, 12.3, 2.5, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 0.7, 3.7, 11.9, 0.4, "INTEGRATIONS SITE / RESEAUX", font_size=14, color=GOLD, bold=True)
integrations = [
    "Feed Instagram en direct sur la page d'accueil",
    "Boutons de partage sur chaque formation et article",
    "Pixel Facebook + TikTok Pixel + GA4 pour tracking",
    "Widget WhatsApp Business pour contact direct",
    "Retargeting publicitaire multi-plateforme",
    "UTM tracking et attribution des conversions",
    "Partage certificats sur LinkedIn automatique",
    "Newsletter integration (Mailchimp / Resend)",
]
for i, item in enumerate(integrations):
    col = i % 2
    row = i // 2
    x = 0.7 + col * 6
    y = 4.15 + row * 0.45
    add_text(sl, x, y, 5.8, 0.4, f"  {item}", font_size=11, color=RGBColor(0xB0, 0xC0, 0xD0))

# ================================================================
# SLIDE 27 - SECTION: PARCOURS UX
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "13", "PARCOURS UTILISATEUR\n& EXPERIENCE (UX)",
              "Le chemin complet de l'apprenant, de la decouverte a l'impact")

# ================================================================
# SLIDE 28 - PARCOURS UX DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "PARCOURS UTILISATEUR", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

ux_steps = [
    ("1", "Decouverte", "Page d'accueil, video inspirante, apercu formations"),
    ("2", "Inscription", "Formulaire simple, SSO Google/Facebook, email"),
    ("3", "Exploration", "Catalogue, filtres, fiches, avis apprenants"),
    ("4", "Paiement", "Checkout securise, multi-options, confirmation"),
    ("5", "Apprentissage", "Videos, quiz, coaching live, defis pratiques"),
    ("6", "Progression", "Dashboard, barre %, badges, rappels auto"),
    ("7", "Certification", "PDF auto + QR code, partage LinkedIn"),
    ("8", "Communaute", "Forum, temoignages, ambassadeur de la vision"),
]
for i, (num, title, desc) in enumerate(ux_steps):
    col = i % 4
    row = i // 4
    x = 0.5 + col * 3.15
    y = 1.15 + row * 2.7
    add_rounded_rect(sl, x, y, 2.95, 2.4, CARD_WHITE)
    # Gold number circle
    add_shape(sl, x + 1.05, y + 0.15, 0.65, 0.65, GOLD, MSO_SHAPE.OVAL)
    add_text(sl, x + 1.05, y + 0.2, 0.65, 0.55, num, font_size=22, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, y + 0.9, 2.95, 0.35, title, font_size=13, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x + 0.1, y + 1.3, 2.75, 0.9, desc, font_size=10, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE - SECTION: CONTACT & DEVIS
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "16", "CONTACT &\nDEVIS PERSONNALIS\u00c9",
              "Nous sommes \u00e0 votre \u00e9coute pour concr\u00e9tiser cette vision ensemble")

# ================================================================
# SLIDE - CONTACT & DEVIS DETAIL
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1.5, -1, 4, GOLD, 0.06), (10, 3.5, 4.5, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "CONTACT & DEVIS PERSONNALIS\u00c9", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

# Left: Contact info
add_rounded_rect(sl, 0.5, 1.2, 6, 5.2, RGBColor(0x14, 0x28, 0x42))
add_shape(sl, 0.5, 1.2, 6, 0.5, GOLD)
add_text(sl, 0.5, 1.23, 6, 0.45, "INFORMATIONS DE CONTACT", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

contact_items = [
    ("\u2709  Email", "contact@messiyagroup.com"),
    ("\u260e  T\u00e9l\u00e9phone", "+33 XX XX XX XX XX"),
    ("\u2302  Si\u00e8ge", "Messiya Group - Division Digitale"),
    ("\u2316  Site Web", "www.messiyagroup.com"),
    ("\u2709  WhatsApp", "+33 XX XX XX XX XX"),
]
for i, (label, value) in enumerate(contact_items):
    y = 1.9 + i * 0.65
    add_text(sl, 0.8, y, 2.5, 0.3, label, font_size=12, color=GOLD, bold=True)
    add_text(sl, 0.8, y + 0.3, 5, 0.3, value, font_size=12, color=RGBColor(0xB0, 0xC0, 0xD0))

add_rounded_rect(sl, 0.7, 5.2, 5.6, 0.9, GOLD)
add_text(sl, 0.7, 5.25, 5.6, 0.4, "PROCHAINE \u00c9TAPE", font_size=12, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
add_text(sl, 0.7, 5.6, 5.6, 0.4, "Planifier un appel de cadrage pour affiner le devis", font_size=11, color=RGBColor(0x1A, 0x2D, 0x50), align=PP_ALIGN.CENTER)

# Right: Devis summary
add_rounded_rect(sl, 6.8, 1.2, 6, 5.2, RGBColor(0x14, 0x28, 0x42))
add_shape(sl, 6.8, 1.2, 6, 0.5, GOLD)
add_text(sl, 6.8, 1.23, 6, 0.45, "CE QUI EST INCLUS DANS LE DEVIS", font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

devis_items = [
    "Maquettes UI/UX compl\u00e8tes (Figma)",
    "D\u00e9veloppement frontend & backend",
    "Int\u00e9gration CRM & paiements",
    "Syst\u00e8me de gestion des formations",
    "Espace membre & dashboard admin",
    "S\u00e9curit\u00e9 renforc\u00e9e (SSL, 2FA, RGPD)",
    "SEO technique & r\u00e9f\u00e9rencement",
    "D\u00e9ploiement & mise en production",
    "Formation de l'\u00e9quipe admin",
    "3 mois de maintenance inclus",
    "Documentation technique compl\u00e8te",
]
for i, item in enumerate(devis_items):
    add_text(sl, 7.1, 1.85 + i * 0.37, 5.5, 0.34, f"\u2713  {item}", font_size=11, color=RGBColor(0xB0, 0xC0, 0xD0))

add_rounded_rect(sl, 7, 5.95, 5.6, 0.4, GOLD)
add_text(sl, 7, 6, 5.6, 0.3, "Devis d\u00e9taill\u00e9 sur demande apr\u00e8s appel de cadrage", font_size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 29 - SECTION: EVOLUTIONS
# ================================================================
sl = prs.slides.add_slide(blank_layout)
section_slide(sl, "17", "\u00c9VOLUTIONS FUTURES\n& ROADMAP",
              "Un plan de livraison structur\u00e9 et des engagements qualit\u00e9 mesurables")

# ================================================================
# SLIDE 30 - EVOLUTIONS DETAIL (like MAG plan de livraison)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1, -1, 3, GOLD, 0.06), (10, 4, 4, GOLD, 0.06)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "PLAN DE LIVRAISON", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)

phases = [
    ("1", "PHASE 1", "Fondation", [
        "Maquettes UI/UX Figma",
        "Charte graphique validee",
        "Setup technique CI/CD",
        "Design system complet",
    ]),
    ("2", "PHASE 2", "Core", [
        "Site vitrine complet",
        "Espace membre & dashboard",
        "Cours enregistres",
        "Paiement en ligne",
    ]),
    ("3", "PHASE 3", "Avancee", [
        "Sessions live Zoom",
        "Forum communautaire",
        "CRM & analytics admin",
        "Certificats auto PDF",
    ]),
    ("4", "PHASE 4", "Lancement", [
        "Tests & audit securite",
        "Deploiement production",
        "Formation equipe admin",
        "SEO & marketing launch",
    ]),
]
for i, (num, phase, subtitle, items_list) in enumerate(phases):
    x = 0.5 + i * 3.15
    add_rounded_rect(sl, x, 1.15, 2.95, 4.2, CARD_WHITE)
    add_shape(sl, x + 1.1, 1.25, 0.65, 0.65, GOLD, MSO_SHAPE.OVAL)
    add_text(sl, x + 1.1, 1.3, 0.65, 0.55, num, font_size=22, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, 2.0, 2.95, 0.35, phase, font_size=14, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl, x, 2.35, 2.95, 0.3, subtitle, font_size=11, color=GOLD, align=PP_ALIGN.CENTER, bold=True)
    add_shape(sl, x + 0.7, 2.7, 1.5, 0.03, GOLD)
    for j, item in enumerate(items_list):
        add_text(sl, x + 0.2, 2.85 + j * 0.5, 2.5, 0.45, f"  {item}", font_size=10, color=TEXT_MUTED)

# Engagements bar
add_rounded_rect(sl, 0.5, 5.6, 12.3, 0.5, GOLD)
add_text(sl, 0.5, 5.65, 12.3, 0.4,
         "ENGAGEMENTS :   Lighthouse > 90   |   < 2s chargement   |   99.9% uptime   |   100% responsive   |   Tous navigateurs",
         font_size=11, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# ================================================================
# SLIDE 31 - VALEUR AJOUTEE MESSIYA GROUP (like MAG slide 26)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [(-1.5, -1, 4, GOLD, 0.06), (10, 3.5, 4.5, GOLD, 0.06), (11, -0.5, 2, GOLD_LIGHT, 0.05)])
top_bar(sl)
footer_bar(sl)

add_text(sl, 0.75, 0.35, 10, 0.6, "VALEUR AJOUTEE MESSIYA GROUP", font_size=28, color=WHITE, bold=True, font_name="Georgia")
add_shape(sl, 0.75, 0.85, 1.5, 0.04, GOLD)
add_text(sl, 0.75, 0.95, 10, 0.4, "Ce que nous apportons en plus pour faire la difference", font_size=13, color=GOLD_LIGHT)

extras = [
    ("PWA (Progressive Web App)", "Application mobile sans telechargement, mode hors-ligne, notifications push pour une experience native sur mobile"),
    ("Intelligence Artificielle", "Recommandations personnalisees de formations, sous-titres auto (Whisper AI), chatbot assistant spirituel 24/7"),
    ("Multi-langue & Accessibilite", "Support Francais, Anglais, Lingala + accessibilite WCAG 2.1 AA pour une portee mondiale"),
    ("Programme de Fidelite", "Points de fidelite, badges gamifies, parrainage avec recompenses automatiques pour les ambassadeurs"),
    ("Tableau de Bord Premium", "Back-office complet : apprenants, formations, paiements, analytics temps reel, rapports de satisfaction"),
    ("Maintenance & Evolution", "Support technique continu, mises a jour securite, evolution fonctionnalites, SLA 99.9% garanti"),
]
for i, (title, desc) in enumerate(extras):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.3
    y = 1.45 + row * 1.6
    add_rounded_rect(sl, x, y, 6, 1.4, RGBColor(0x14, 0x28, 0x42))
    add_shape(sl, x, y, 0.08, 1.4, GOLD)
    add_text(sl, x + 0.25, y + 0.08, 0.35, 0.35, "+", font_size=20, color=GOLD, bold=True)
    add_text(sl, x + 0.65, y + 0.1, 5.1, 0.35, title, font_size=13, color=WHITE, bold=True)
    add_text(sl, x + 0.65, y + 0.5, 5.1, 0.8, desc, font_size=10, color=RGBColor(0xB0, 0xC0, 0xD0))

# ================================================================
# SLIDE 32 - MERCI (like MAG slide 27)
# ================================================================
sl = prs.slides.add_slide(blank_layout)
set_bg(sl, BG)
decor_circles(sl, [
    (-2, -2, 6, GOLD, 0.08),
    (-0.7, -1, 4, GOLD, 0.06),
    (8, 3, 5.5, GOLD, 0.06),
    (9.5, 4, 3.5, GOLD, 0.08),
    (10.5, -0.8, 2.8, GOLD_LIGHT, 0.05),
])
top_bar(sl)

sl.shapes.add_picture(LOGO, Inches(4.5), Inches(0.3), Inches(4.3), Inches(3.1))

add_text(sl, 1, 3.5, 11, 1.2,
         "MERCI POUR\nVOTRE CONFIANCE",
         font_size=40, color=WHITE, bold=True, font_name="Georgia", align=PP_ALIGN.CENTER)

add_shape(sl, 5, 4.7, 3.3, 0.04, GOLD)

add_text(sl, 1, 4.9, 11, 0.5,
         "Lord Lombo Academie x Messiya Group",
         font_size=18, color=GOLD, align=PP_ALIGN.CENTER, font_name="Georgia")

add_text(sl, 1.5, 5.5, 10, 0.7,
         "Ensemble, faisons de Lord Lombo Academie la reference\n"
         "de la formation spirituelle et du leadership dans l'espace francophone.",
         font_size=14, color=RGBColor(0xB0, 0xC0, 0xD0), align=PP_ALIGN.CENTER)

add_rounded_rect(sl, 3.5, 6.3, 6, 0.5, RGBColor(0x14, 0x28, 0x42))
add_text(sl, 3.5, 6.35, 6, 0.4,
         "Contact : Messiya Group\nDivision Digitale",
         font_size=9, color=RGBColor(0x88, 0x99, 0xAA), align=PP_ALIGN.CENTER)

add_shape(sl, 0, H_IN - 0.4, W_IN, 0.4, FOOTER_BG)
add_text(sl, 2, H_IN - 0.38, 9, 0.35,
         "CONFIDENTIEL  |  Fevrier 2026  |  Version 1.0",
         font_size=10, color=WHITE, align=PP_ALIGN.CENTER, bold=True)

# SAVE
output = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Lord_Lombo_Academie_V3.pptx")
prs.save(output)
print(f"Saved: {output} ({len(prs.slides)} slides)")
